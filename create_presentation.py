import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_investor_ppt():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_TERRACOTTA = RGBColor(224, 83, 38)
    COLOR_CHARCOAL = RGBColor(30, 27, 24)
    COLOR_CREAM = RGBColor(250, 247, 242)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_GOLD = RGBColor(245, 158, 11)
    COLOR_EMERALD = RGBColor(4, 120, 87)
    COLOR_SLATE_GRAY = RGBColor(100, 116, 139)
    COLOR_LIGHT_CARD = RGBColor(245, 242, 235)

    def add_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="KOOKINDIA INVESTOR PRESENTATION"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_TERRACOTTA

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_CHARCOAL

    # ================= SLIDE 1: TITLE SLIDE =================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide1, COLOR_CHARCOAL)

    # Title box
    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.5))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "KookIndia"
    p0.font.size = Pt(54)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_TERRACOTTA
    p0.space_after = Pt(10)

    p1 = tf1.add_paragraph()
    p1.text = "Empowering 300M+ Indian Homemakers to Deliver Authentic Regional Home Food."
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(20)

    p2 = tf1.add_paragraph()
    p2.text = "Seed Round Investment Pitch Deck & Product Blueprint | Non-Technical Investor Guide"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_GOLD

    # Footer Box
    fbox = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.0))
    tff = fbox.text_frame
    pf = tff.paragraphs[0]
    pf.text = "Presenter: Nithin Joseph & Engineering Team  |  GitHub: https://github.com/njoseph11224-tech/kookindia.git"
    pf.font.size = Pt(11)
    pf.font.color.rgb = COLOR_SLATE_GRAY

    # ================= SLIDE 2: THE BIG PROBLEM =================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide2, COLOR_CREAM)
    add_header(slide2, "The Market Problem: What Breaks in Food Delivery Today?")

    # Problem Card 1: Customer Pain
    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_WHITE
    c1.line.color.rgb = COLOR_SLATE_GRAY

    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "1. Customer Pain: Restaurant Fatigue"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TERRACOTTA
    p.space_after = Pt(14)

    items1 = [
        "85 Million+ urban professionals & students in Indian metros eat out daily.",
        "Commercial dark kitchens overuse heavy palm oil, artificial colors & spices.",
        "High daily cost: Restaurant meals are expensive (₹350-₹500 per meal).",
        "Zero authentic regional options: Hard to find genuine Gujarati, Bengali, Chettinad, or Jain home food."
    ]
    for item in items1:
        pi = tf_c1.add_paragraph()
        pi.text = "• " + item
        pi.font.size = Pt(12)
        pi.font.color.rgb = COLOR_CHARCOAL
        pi.space_after = Pt(10)

    # Problem Card 2: Homemaker Pain
    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.4), Inches(4.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_WHITE
    c2.line.color.rgb = COLOR_SLATE_GRAY

    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "2. Homemaker Pain: Zero Monetization"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TERRACOTTA
    p.space_after = Pt(14)

    items2 = [
        "300 Million+ Indian homemakers possess elite culinary skills.",
        "High barrier to entry: Opening a cloud kitchen costs ₹15-20 Lakhs in upfront capital.",
        "Heavy platform commissions: Zomato & Swiggy charge 25% - 30% from kitchens.",
        "Marketing & delivery hurdle: Homemakers cannot manage delivery drivers or online marketing alone."
    ]
    for item in items2:
        pi = tf_c2.add_paragraph()
        pi.text = "• " + item
        pi.font.size = Pt(12)
        pi.font.color.rgb = COLOR_CHARCOAL
        pi.space_after = Pt(10)

    # ================= SLIDE 3: THE SOLUTION =================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide3, COLOR_CREAM)
    add_header(slide3, "The KookIndia Solution: Uber for Home Chefs")

    sol_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True

    p = tf_sol.paragraphs[0]
    p.text = "A Tech Marketplace Connecting Verified Home Kitchens with Urban Food Lovers."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CHARCOAL
    p.space_after = Pt(20)

    points = [
        ("🥗 Fresh Ghar Ka Khana", "Cooked in small batches by homemakers using cold-pressed oils, fresh vegetables, and hand-ground spices."),
        ("🛵 Zero Delivery Logistics Overhead", "We integrate directly with Dunzo, Porter & Borzo to pick up meals straight from the cook's home."),
        ("⚡ Pre-Order & Ready-Batch Model", "Homemakers cook with 2-hour advance notices or prepare ready-to-serve batches, eliminating food waste."),
        ("🛡️ 100% FSSAI Hygiene Verified", "We assist home cooks in obtaining their ₹100/yr FSSAI license and conduct virtual kitchen inspections."),
        ("📅 Daily Office Tiffin Subscriptions", "Recurring monthly meal plans for bachelors and corporate workers, providing steady revenue.")
    ]

    for title, desc in points:
        pi = tf_sol.add_paragraph()
        pi.text = f"{title}: {desc}"
        pi.font.size = Pt(13)
        pi.font.color.rgb = COLOR_CHARCOAL
        pi.space_after = Pt(12)

    # ================= SLIDE 4: HOW IT WORKS =================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide4, COLOR_CREAM)
    add_header(slide4, "How KookIndia Works (3 Simple Steps)")

    step_data = [
        ("Step 1: Discover Local Kitchens", "Customer selects city & locality (e.g. Bangalore - Indiranagar). Filters by North Indian, Chettinad, Gujarati, Bengali, or Jain thalis.", COLOR_TERRACOTTA),
        ("Step 2: Pre-Order or Subscribe", "Customer picks a single meal or subscribes to a 30-day Mon-Fri lunch plan. Items can be combined from multiple chefs in 1 basket.", COLOR_CHARCOAL),
        ("Step 3: Fresh Doorstep Delivery", "The home cook receives the order, cooks it fresh, and hyperlocal drivers (Dunzo/Porter) deliver piping hot to the customer.", COLOR_EMERALD)
    ]

    for idx, (title, desc, color) in enumerate(step_data):
        left_pos = Inches(1.0 + idx * 3.9)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.2), Inches(3.6), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.space_after = Pt(14)

        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_CHARCOAL

    # ================= SLIDE 5: MARKET OPPORTUNITY =================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide5, COLOR_CREAM)
    add_header(slide5, "Market Opportunity: $40 Billion+ Indian Food Industry")

    metrics = [
        ("$40 Billion+", "Total Addressable Market (TAM)", "Indian Food Services & Delivery Industry by 2026", COLOR_TERRACOTTA),
        ("$6 Billion+", "Serviceable Market (SAM)", "Home Tiffin & Regional Meal Market in Top 10 Indian Cities", COLOR_GOLD),
        ("$150 Million", "3-Year SOM Goal", "Gross Merchandise Value target across 6 launch metros", COLOR_EMERALD)
    ]

    for idx, (stat, title, desc, color) in enumerate(metrics):
        left_pos = Inches(1.0 + idx * 3.9)
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.2), Inches(3.6), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_WHITE
        box.line.color.rgb = COLOR_LIGHT_CARD

        tf = box.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = stat
        p0.font.size = Pt(36)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.space_after = Pt(10)

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CHARCOAL
        p1.space_after = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_SLATE_GRAY

    # ================= SLIDE 6: BUSINESS MODEL =================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide6, COLOR_CREAM)
    add_header(slide6, "Business Model: How KookIndia Generates Revenue")

    rev_streams = [
        ("1. Marketplace Commission (12% - 15%)", "We charge home cooks a small 12-15% fee per completed dish order (compared to Zomato's heavy 25-30% commission)."),
        ("2. Monthly Tiffin Subscriptions", "Recurring monthly upfront subscription payments from office goers (High retention & zero customer churn)."),
        ("3. Hyperlocal Delivery Markup", "Nominal delivery fee passed to logistics partners (Dunzo/Porter) with platform margin."),
        ("4. Featured Kitchen Ads", "Promoted banner listings for home cooks wishing to boost order volumes on top of city feeds.")
    ]

    for idx, (title, desc) in enumerate(rev_streams):
        row = idx // 2
        col = idx % 2
        top_pos = Inches(2.0 + row * 2.5)
        left_pos = Inches(1.0 + col * 5.8)

        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.5), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_LIGHT_CARD

        tf = card.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TERRACOTTA
        p0.space_after = Pt(8)

        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_CHARCOAL

    # ================= SLIDE 7: UNIT ECONOMICS =================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide7, COLOR_CREAM)
    add_header(slide7, "Unit Economics: Exceptional 20x LTV:CAC Ratio")

    # Table of Economics
    table_shape = slide7.shapes.add_table(7, 3, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    table = table_shape.table

    headers = ["Metric", "Amount (₹)", "Plain English Explanation"]
    data = [
        ["Average Order Value (AOV)", "₹380.00", "Typical customer basket size (2 regional meals or thali)"],
        ["Platform Commission (14%)", "₹53.20", "Our core marketplace cut"],
        ["Platform Hygiene Fee", "₹15.00", "Fixed per-order platform fee"],
        ["Gross Margin / Order", "₹68.20 (18%)", "Net revenue generated per order after payment fees"],
        ["Customer Acquisition Cost (CAC)", "₹120.00", "Cost to acquire 1 new active customer via digital ads"],
        ["Customer Lifetime Value (LTV)", "₹2,400.00", "Total revenue generated per customer (20 orders/year)"]
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CHARCOAL
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 == 0 else COLOR_LIGHT_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_CHARCOAL
            if c == 1:
                p.font.bold = True

    # ================= SLIDE 8: TECH & PROTOTYPE =================
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide8, COLOR_CREAM)
    add_header(slide8, "Live Working Prototype & Technology Architecture")

    proto_features = [
        ("🛒 Multi-Chef Grouped Cart", "Customers can add dishes from multiple neighborhood cooks in 1 order. Our backend splits the order and dispatches drivers to both kitchens in parallel."),
        ("👨‍🍳 Homemaker Kitchen Portal", "Simple dashboard for home cooks: Kitchen OPEN/CLOSED toggle, Photo Upload file picker with live preview, Advance Pre-Batch setup, and 1-click Bulk Order Processing."),
        ("🛡️ Admin Verification & Support", "Operations portal for FSSAI license verification (`21223190004512`), virtual kitchen photo audits, and 15-minute customer complaint resolution."),
        ("📱 Modern Tech Stack", "Built with Next.js 16 (App Router), TypeScript, Tailwind CSS, and Supabase PostgreSQL.")
    ]

    for idx, (title, desc) in enumerate(proto_features):
        row = idx // 2
        col = idx % 2
        top_pos = Inches(2.0 + row * 2.5)
        left_pos = Inches(1.0 + col * 5.8)

        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.5), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_LIGHT_CARD

        tf = card.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TERRACOTTA
        p0.space_after = Pt(8)

        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_CHARCOAL

    # ================= SLIDE 9: COMPETITIVE ADVANTAGE =================
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide9, COLOR_CREAM)
    add_header(slide9, "Competitive Advantage: Why We Win")

    comp_table_shape = slide9.shapes.add_table(6, 4, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    comp_table = comp_table_shape.table

    comp_headers = ["Feature", "KookIndia", "Zomato / Swiggy", "Unorganized Tiffin Vendors"]
    comp_data = [
        ["Food Supply", "Verified Homemakers", "Commercial Restaurants", "Unorganized Local Cooks"],
        ["Taste & Hygiene", "Homestyle, Low-Oil", "Commercial Heavy Oil", "Inconsistent"],
        ["Commission Charged", "12% - 15%", "25% - 30%", "0%"],
        ["Multi-Chef Basket", "YES (Built-in)", "NO", "NO"],
        ["Tiffin Subscriptions", "YES (Built-in)", "Limited", "Cash Only"]
    ]

    for c, h in enumerate(comp_headers):
        cell = comp_table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CHARCOAL
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for r, row in enumerate(comp_data):
        for c, val in enumerate(row):
            cell = comp_table.cell(r + 1, c)
            cell.fill.solid()
            if c == 1:
                cell.fill.fore_color.rgb = RGBColor(254, 243, 199) # Highlight KookIndia column
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 == 0 else COLOR_LIGHT_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TERRACOTTA if c == 1 else COLOR_CHARCOAL
            if c == 1:
                p.font.bold = True

    # ================= SLIDE 10: ROADMAP =================
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide10, COLOR_CREAM)
    add_header(slide10, "12-Month Execution Roadmap")

    phases = [
        ("Phase 1: Soft Launch (Q3 2026)", "Launch in Bangalore (Indiranagar, Koramangala) & Gurgaon (DLF Phase 1-5). Onboard 500 verified home kitchens & launch web app.", COLOR_TERRACOTTA),
        ("Phase 2: Regional Scale (Q4 2026)", "Expand to Mumbai (Ghatkopar, Powai) & Hyderabad (Banjara Hills). Release iOS & Android React Native Mobile Apps.", COLOR_GOLD),
        ("Phase 3: Hyper-Growth (Q1 2027)", "Cross 50,000 monthly orders. Launch B2B Corporate Tiffin Subscriptions & expand to Pune and Delhi NCR.", COLOR_EMERALD)
    ]

    for idx, (title, desc, color) in enumerate(phases):
        top_pos = Inches(2.0 + idx * 1.6)
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), top_pos, Inches(11.3), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.space_after = Pt(4)

        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_CHARCOAL

    # ================= SLIDE 11: INVESTMENT ASK =================
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide11, COLOR_CREAM)
    add_header(slide11, "Investment Ask & Allocation of Funds")

    # Large Ask Box
    ask_box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.2))
    ask_box.fill.solid()
    ask_box.fill.fore_color.rgb = COLOR_CHARCOAL

    tf_ask = ask_box.text_frame
    tf_ask.word_wrap = True
    p = tf_ask.paragraphs[0]
    p.text = "Seeking $500,000 (Seed Round / ~₹4.1 Crore)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.alignment = PP_ALIGN.CENTER

    use_of_funds = [
        ("40% - Homemaker Onboarding & FSSAI Support", "₹1.64 Crore allocated to kitchen onboarding, virtual audits, and licensing assistance."),
        ("35% - Customer Acquisition & Ads", "₹1.43 Crore allocated to digital ads, influencer marketing, and customer referral bonuses."),
        ("15% - Mobile App Development", "₹61.5 Lakhs allocated to iOS & Android native apps."),
        ("10% - Operations & Audit Compliance", "₹41 Lakhs allocated to customer support and quality assurance.")
    ]

    for idx, (title, desc) in enumerate(use_of_funds):
        row = idx // 2
        col = idx % 2
        top_pos = Inches(3.3 + row * 1.9)
        left_pos = Inches(1.0 + col * 5.8)

        card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.5), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_LIGHT_CARD

        tf = card.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TERRACOTTA
        p0.space_after = Pt(4)

        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(10.5)
        p1.font.color.rgb = COLOR_CHARCOAL

    # ================= SLIDE 12: CONTACT =================
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide12, COLOR_CHARCOAL)

    cbox = slide12.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.0))
    tfc = cbox.text_frame
    tfc.word_wrap = True

    p0 = tfc.paragraphs[0]
    p0.text = "Join Us in Revolutionizing Home Food in India"
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_TERRACOTTA
    p0.space_after = Pt(20)

    contacts = [
        "🏢 Company: KookIndia Technologies Pvt Ltd",
        "✉️ Investor Email: investor@kookindia.com",
        "🌐 Live Prototype URL: http://localhost:3000",
        "📂 GitHub Repository: https://github.com/njoseph11224-tech/kookindia.git",
        "👤 Lead Engineering: Nithin Joseph & Engineering Team"
    ]

    for c in contacts:
        p = tfc.add_paragraph()
        p.text = c
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(12)

    output_path = r"C:\Nithin-Learning\Projects\kookindia\KookIndia_Investor_Presentation.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Created PowerPoint Presentation at {output_path}")

if __name__ == "__main__":
    build_investor_ppt()
